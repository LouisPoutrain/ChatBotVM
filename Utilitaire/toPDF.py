#!/usr/bin/env python3

import subprocess
import sys
from pathlib import Path
import shutil

# Dépendances optionnelles:
# pip install pillow reportlab
try:
    from PIL import Image
except Exception:
    Image = None

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
except Exception:
    canvas = None
    A4 = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DOWNLOADS_DIR = PROJECT_ROOT / "data" / "downloads"
OUTPUT_DIR = DOWNLOADS_DIR / "pdf"
IS_MACOS = sys.platform == "darwin"


IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".webp"}
TEXT_EXTS = {".txt", ".md", ".csv", ".log", ".json"}
OFFICE_EXTS = {".doc", ".docx", ".odt", ".rtf", ".xls", ".xlsx", ".xlsm", ".ods", ".ppt", ".pptx", ".odp"}


def unique_pdf_path(base_name: str, out_dir: Path) -> Path:
    target = out_dir / f"{base_name}.pdf"
    i = 1
    while target.exists():
        target = out_dir / f"{base_name}_{i}.pdf"
        i += 1
    return target


def output_pdf_path(src: Path, output_root: Path) -> Path:
    # Enregistrer directement dans le dossier de sortie plat (sans recréer l'arborescence)
    output_root.mkdir(parents=True, exist_ok=True)
    return output_root / f"{src.stem}.pdf"


def escape_applescript_string(value: str) -> str:
    return value.replace("\\", "\\\\").replace('"', '\\"')


def run_osascript(script: str) -> None:
    completed = subprocess.run(
        ["osascript", "-e", script],
        capture_output=True,
        text=True,
        check=False,
    )
    if completed.returncode != 0:
        stderr = completed.stderr.strip() or completed.stdout.strip() or "Erreur inconnue"
        raise RuntimeError(stderr)


def convert_docx_to_pdf_applescript(src: Path, dst: Path) -> None:
    if not IS_MACOS:
        raise RuntimeError("AppleScript est disponible uniquement sur macOS")

    input_path = escape_applescript_string(str(src))
    output_path = escape_applescript_string(str(dst))
    script = f'''
set inputPath to POSIX file "{input_path}"
set outputPath to "{output_path}"
tell application "Microsoft Word"
    activate
    open inputPath
    delay 1
    set theDoc to active document
    save as theDoc file name outputPath file format format PDF
    close theDoc saving no
end tell
'''
    run_osascript(script)


def convert_xlsx_to_pdf_applescript(src: Path, dst: Path) -> None:
    if not IS_MACOS:
        raise RuntimeError("AppleScript est disponible uniquement sur macOS")

    input_path = escape_applescript_string(str(src))
    output_path = escape_applescript_string(str(dst))
    script = f'''
set inputPath to POSIX file "{input_path}"
set outputPath to POSIX file "{output_path}"
tell application "Microsoft Excel"
    activate
    open inputPath
    delay 1
    set theWorkbook to active workbook
    save workbook as theWorkbook filename outputPath file format PDF file format
    close theWorkbook saving no
end tell
'''
    run_osascript(script)


def convert_image_to_pdf(src: Path, dst: Path) -> None:
    if Image is None:
        raise RuntimeError("Pillow non installé (pip install pillow)")
    with Image.open(src) as img:
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img.save(dst, "PDF", resolution=100.0)


def convert_text_to_pdf(src: Path, dst: Path) -> None:
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 13

    with src.open("r", encoding="utf-8", errors="replace") as f:
        for line in f:
            line = line.rstrip("\n")
            # Découpe simple des lignes trop longues
            chunk_size = 100
            parts = [line[i:i + chunk_size] for i in range(0, len(line), chunk_size)] or [""]
            for part in parts:
                if y < margin:
                    c.showPage()
                    y = height - margin
                c.drawString(margin, y, part)
                y -= line_height

    c.save()


def convert_docx_to_pdf_text(src: Path, dst: Path) -> None:
    if Document is None:
        raise RuntimeError("python-docx non installé (pip install python-docx)")
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")

    document = Document(str(src))
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 13

    def write_line(text: str) -> None:
        nonlocal y
        chunk_size = 100
        parts = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)] or [""]
        for part in parts:
            if y < margin:
                c.showPage()
                y = height - margin
            c.drawString(margin, y, part)
            y -= line_height

    for paragraph in document.paragraphs:
        write_line(paragraph.text)

    for table in document.tables:
        write_line("")
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            write_line(row_text)

    c.save()


def convert_pptx_to_pdf(src: Path, dst: Path) -> None:
    if Presentation is None:
        raise RuntimeError("python-pptx non installé (pip install python-pptx)")
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")

    presentation = Presentation(str(src))
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 14

    def write_line(text: str) -> None:
        nonlocal y
        chunk_size = 95
        parts = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)] or [""]
        for part in parts:
            if y < margin:
                c.showPage()
                y = height - margin
            c.drawString(margin, y, part)
            y -= line_height

    for slide_index, slide in enumerate(presentation.slides, start=1):
        write_line(f"Slide {slide_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    write_line(line)
        write_line("")

    c.save()


def convert_xlsx_to_pdf_text(src: Path, dst: Path) -> None:
    if load_workbook is None:
        raise RuntimeError("openpyxl non installé (pip install openpyxl)")
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")

    workbook = load_workbook(str(src), data_only=True)
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 13

    def write_line(text: str) -> None:
        nonlocal y
        chunk_size = 100
        parts = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)] or [""]
        for part in parts:
            if y < margin:
                c.showPage()
                y = height - margin
            c.drawString(margin, y, part)
            y -= line_height

    for sheet in workbook.worksheets:
        write_line(f"Feuille: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join("" if cell is None else str(cell) for cell in row)
            write_line(row_text)
        write_line("")

    c.save()


def convert_docx_to_pdf(src: Path, dst: Path) -> None:
    if IS_MACOS:
        try:
            convert_docx_to_pdf_applescript(src, dst)
            return
        except Exception as exc:
            print(f"[AVERTISSEMENT] AppleScript Word a échoué pour {src.name}: {exc}. Repli sur le mode texte.")

    convert_docx_to_pdf_text(src, dst)


def convert_xlsx_to_pdf(src: Path, dst: Path) -> None:
    if IS_MACOS:
        try:
            convert_xlsx_to_pdf_applescript(src, dst)
            return
        except Exception as exc:
            print(f"[AVERTISSEMENT] AppleScript Excel a échoué pour {src.name}: {exc}. Repli sur le mode texte.")

    convert_xlsx_to_pdf_text(src, dst)


def convert_ods_to_pdf(src: Path, dst: Path) -> None:
    raise RuntimeError("Le format .ods n'est pas converti automatiquement sans LibreOffice")


def convert_office_to_pdf(src: Path, dst: Path) -> bool:
    ext = src.suffix.lower()

    try:
        if ext == ".docx":
            convert_docx_to_pdf(src, dst)
            return True
        if ext == ".pptx":
            convert_pptx_to_pdf(src, dst)
            return True
        if ext in {".xlsx", ".xlsm"}:
            convert_xlsx_to_pdf(src, dst)
            return True
        if ext == ".ods":
            convert_ods_to_pdf(src, dst)
            return True
    except Exception as exc:
        print(f"[ERREUR] {src.name}: {exc}")
        return False

    print(f"[IGNORÉ] {src.name}: format Office non pris en charge sans LibreOffice")
    return False


def main():
    """
    Outil de conversion universelle de documents vers PDF.
    
    Supporte:
    - Microsoft Office (DOCX, XLSX, PPTX) via AppleScript (sur Mac) ou LibreOffice/Extraction texte.
    - OpenDocument (ODS).
    - Images (JPG, PNG).
    - Texte brut (.txt).
    
    Parcourt un répertoire source et génère des PDF dans un répertoire cible.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    converted = 0
    skipped = 0
    errors = 0

    for f in DOWNLOADS_DIR.rglob("*"):
        if not f.is_file():
            continue
        if f.parent == OUTPUT_DIR:
            continue
        if OUTPUT_DIR in f.parents:
            continue
        if f.suffix.lower() == ".pdf":
            continue

        ext = f.suffix.lower()
        dst = output_pdf_path(f, OUTPUT_DIR)

        # Si le PDF de destination existe déjà, on l'ignore (skip)
        if dst.exists():
            skipped += 1
            continue

        try:
            if ext == ".pdf":
                shutil.copy2(f, dst)
                converted += 1

            elif ext in IMAGE_EXTS:
                convert_image_to_pdf(f, dst)
                converted += 1

            elif ext in TEXT_EXTS:
                convert_text_to_pdf(f, dst)
                converted += 1

            elif ext in OFFICE_EXTS:
                if convert_office_to_pdf(f, dst):
                    converted += 1
                else:
                    skipped += 1

            else:
                skipped += 1

        except Exception as e:
            errors += 1
            print(f"[ERREUR] {f.relative_to(DOWNLOADS_DIR)}: {e}")

    print(f"Terminé. Convertis: {converted}, Ignorés: {skipped}, Erreurs: {errors}")
    print(f"PDF disponibles dans: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()