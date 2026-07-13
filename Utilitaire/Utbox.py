from __future__ import annotations

import json
import re
import shutil
import zipfile
import subprocess
import sys
from pathlib import Path
from urllib.parse import quote, urlparse
from xml.etree import ElementTree as ET
import requests
from loguru import logger

PROJECT_ROOT = Path(__file__).resolve().parent.parent

def setup_global_logger(log_dir_name: str = "Log", log_prefix: str = ""):
    """
    Configure loguru to use a centralized log directory with daily rotation
    and month-based subdirectories.
    """
    logger.remove()  # Remove default console logger
    logger.add(sys.stderr, level="INFO")
    
    log_dir = PROJECT_ROOT / log_dir_name
    
    filename = "{time:YYYY-MM-DD}"
    if log_prefix:
        filename += f"_{log_prefix}"
    filename += ".log"
    
    log_pattern = str(log_dir / "{time:YYYY-MM}" / filename)
    
    logger.add(
        log_pattern,
        rotation="1 day",
        retention="1 year",
        level="DEBUG",
        encoding="utf-8",
        enqueue=True, # Thread-safe
        backtrace=True,
        diagnose=True
    )
    return logger
# Dépendances optionnelles pour la conversion PDF
try:
    from PIL import Image
except Exception:
    Image = None

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
except Exception:
    canvas = None
    A4 = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

IS_MACOS = sys.platform == "darwin"
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".webp"}
TEXT_EXTS = {".txt", ".md", ".csv", ".log", ".json"}
OFFICE_EXTS = {".doc", ".docx", ".odt", ".rtf", ".xls", ".xlsx", ".xlsm", ".ods", ".ppt", ".pptx", ".odp"}


def unique_pdf_path(base_name: str, out_dir: Path) -> Path:
    target = out_dir / f"{base_name}.pdf"
    i = 1
    while target.exists():
        target = out_dir / f"{base_name}_{i}.pdf"
        i += 1
    return target


def escape_applescript_string(value: str) -> str:
    return value.replace("\\", "\\\\").replace('"', '\\"')


def run_osascript(script: str) -> None:
    completed = subprocess.run(
        ["osascript", "-e", script],
        capture_output=True,
        text=True,
        check=False,
    )
    if completed.returncode != 0:
        stderr = completed.stderr.strip() or completed.stdout.strip() or "Erreur inconnue"
        raise RuntimeError(stderr)


def convert_docx_to_pdf_applescript(src: Path, dst: Path) -> None:
    if not IS_MACOS:
        raise RuntimeError("AppleScript est disponible uniquement sur macOS")

    input_path = escape_applescript_string(str(src))
    output_path = escape_applescript_string(str(dst))
    script = f'''
set inputPath to POSIX file "{input_path}"
set outputPath to "{output_path}"
tell application "Microsoft Word"
    activate
    open inputPath
    delay 1
    set theDoc to active document
    save as theDoc file name outputPath file format format PDF
    close theDoc saving no
end tell
'''
    run_osascript(script)


def convert_xlsx_to_pdf_applescript(src: Path, dst: Path) -> None:
    if not IS_MACOS:
        raise RuntimeError("AppleScript est disponible uniquement sur macOS")

    input_path = escape_applescript_string(str(src))
    output_path = escape_applescript_string(str(dst))
    script = f'''
set inputPath to POSIX file "{input_path}"
set outputPath to POSIX file "{output_path}"
tell application "Microsoft Excel"
    activate
    open inputPath
    delay 1
    set theWorkbook to active workbook
    save workbook as theWorkbook filename outputPath file format PDF file format
    close theWorkbook saving no
end tell
'''
    run_osascript(script)


def convert_image_to_pdf(src: Path, dst: Path) -> None:
    if Image is None:
        raise RuntimeError("Pillow non installé (pip install pillow)")
    with Image.open(src) as img:
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img.save(dst, "PDF", resolution=100.0)


def convert_text_to_pdf(src: Path, dst: Path) -> None:
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 13

    with src.open("r", encoding="utf-8", errors="replace") as f:
        for line in f:
            line = line.rstrip("\n")
            chunk_size = 100
            parts = [line[i:i + chunk_size] for i in range(0, len(line), chunk_size)] or [""]
            for part in parts:
                if y < margin:
                    c.showPage()
                    y = height - margin
                c.drawString(margin, y, part)
                y -= line_height

    c.save()


def convert_docx_to_pdf_text(src: Path, dst: Path) -> None:
    if Document is None:
        raise RuntimeError("python-docx non installé (pip install python-docx)")
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")

    document = Document(str(src))
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 13

    def write_line(text: str) -> None:
        nonlocal y
        chunk_size = 100
        parts = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)] or [""]
        for part in parts:
            if y < margin:
                c.showPage()
                y = height - margin
            c.drawString(margin, y, part)
            y -= line_height

    for paragraph in document.paragraphs:
        write_line(paragraph.text)

    for table in document.tables:
        write_line("")
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            write_line(row_text)

    c.save()


def convert_pptx_to_pdf(src: Path, dst: Path) -> None:
    if Presentation is None:
        raise RuntimeError("python-pptx non installé (pip install python-pptx)")
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")

    presentation = Presentation(str(src))
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 14

    def write_line(text: str) -> None:
        nonlocal y
        chunk_size = 95
        parts = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)] or [""]
        for part in parts:
            if y < margin:
                c.showPage()
                y = height - margin
            c.drawString(margin, y, part)
            y -= line_height

    for slide_index, slide in enumerate(presentation.slides, start=1):
        write_line(f"Slide {slide_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    write_line(line)
        write_line("")

    c.save()


def convert_xlsx_to_pdf_text(src: Path, dst: Path) -> None:
    if load_workbook is None:
        raise RuntimeError("openpyxl non installé (pip install openpyxl)")
    if canvas is None:
        raise RuntimeError("reportlab non installé (pip install reportlab)")

    workbook = load_workbook(str(src), data_only=True)
    c = canvas.Canvas(str(dst), pagesize=A4)
    width, height = A4
    margin = 40
    y = height - margin
    line_height = 13

    def write_line(text: str) -> None:
        nonlocal y
        chunk_size = 100
        parts = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)] or [""]
        for part in parts:
            if y < margin:
                c.showPage()
                y = height - margin
            c.drawString(margin, y, part)
            y -= line_height

    for sheet in workbook.worksheets:
        write_line(f"Feuille: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join("" if cell is None else str(cell) for cell in row)
            write_line(row_text)
        write_line("")

    c.save()


def convert_docx_to_pdf(src: Path, dst: Path) -> None:
    if IS_MACOS:
        try:
            convert_docx_to_pdf_applescript(src, dst)
            return
        except Exception as exc:
            print(f"[AVERTISSEMENT] AppleScript Word a échoué pour {src.name}: {exc}. Repli sur le mode texte.")

    convert_docx_to_pdf_text(src, dst)


def convert_xlsx_to_pdf(src: Path, dst: Path) -> None:
    if IS_MACOS:
        try:
            convert_xlsx_to_pdf_applescript(src, dst)
            return
        except Exception as exc:
            print(f"[AVERTISSEMENT] AppleScript Excel a échoué pour {src.name}: {exc}. Repli sur le mode texte.")

    convert_xlsx_to_pdf_text(src, dst)


def convert_ods_to_pdf(src: Path, dst: Path) -> None:
    raise RuntimeError("Le format .ods n'est pas converti automatiquement sans LibreOffice")


def convert_office_to_pdf(src: Path, dst: Path) -> bool:
    ext = src.suffix.lower()

    try:
        if ext == ".docx":
            convert_docx_to_pdf(src, dst)
            return True
        if ext == ".pptx":
            convert_pptx_to_pdf(src, dst)
            return True
        if ext in {".xlsx", ".xlsm"}:
            convert_xlsx_to_pdf(src, dst)
            return True
        if ext == ".ods":
            convert_ods_to_pdf(src, dst)
            return True
    except Exception as exc:
        print(f"[ERREUR] {src.name}: {exc}")
        return False

    print(f"[IGNORÉ] {src.name}: format Office non pris en charge sans LibreOffice")
    return False


class UTBoxDownloader:
    """
    Classe permettant de télécharger et de synchroniser des fichiers depuis UTBOX (Nextcloud)
    via l'API WebDAV ou en téléchargeant directement une archive ZIP.
    Gère également la conversion des fichiers Office en PDF si nécessaire.
    """
    def __init__(self, timeout: int = 60):
        self.timeout = timeout
        self.session = requests.Session()
        self.session.headers.update(
            {
                "User-Agent": "Mozilla/5.0 (UTBoxDownloader/1.0)",
                "Accept": "*/*",
            }
        )

    @staticmethod
    def _extract_share_info(url: str) -> tuple[str, str]:
        parsed = urlparse(url)
        match = re.search(r"/s/([^/]+)", parsed.path)
        if not match:
            raise ValueError("URL UTBox invalide: token de partage introuvable.")

        base_url = f"{parsed.scheme}://{parsed.netloc}"
        share_token = match.group(1)
        return base_url, share_token

    def _resolve_webdav_root(self, base_url: str, share_token: str) -> str:
        candidates = [
            f"{base_url}/public.php/dav/files/{share_token}/",
            f"{base_url}/public.php/webdav/",
        ]

        for candidate in candidates:
            try:
                resp = self.session.request(
                    method="PROPFIND",
                    url=candidate,
                    headers={"Depth": "0"},
                    auth=(share_token, ""),
                    timeout=self.timeout,
                )
                if resp.status_code in (200, 207):
                    return candidate
            except requests.RequestException:
                continue

        raise requests.HTTPError("Impossible de trouver un endpoint WebDAV public valide pour ce lien UTBox.")

    @staticmethod
    def _filename_from_headers(headers: dict) -> str | None:
        content_disposition = headers.get("Content-Disposition", "")
        match = re.search(r'filename\*?=(?:UTF-8\'\')?"?([^";]+)"?', content_disposition)
        if match:
            return match.group(1).strip()
        return None

    @staticmethod
    def _filename_from_url(url: str) -> str:
        path = urlparse(url).path
        name = Path(path).name
        return name if name and name != "download" else "utbox_download.bin"

    @staticmethod
    def _unique_path(path: Path) -> Path:
        if not path.exists():
            return path

        index = 1
        while True:
            candidate = path.with_name(f"{path.stem}_{index}{path.suffix}")
            if not candidate.exists():
                return candidate
            index += 1

    @staticmethod
    def _copy_stream(src, dst, buffer_size: int = 8 * 1024 * 1024) -> None:
        shutil.copyfileobj(src, dst, length=buffer_size)

    @staticmethod
    def _propfind_parse(xml_text: str) -> list[tuple[str, bool, str]]:
        ns = {
            "d": "DAV:",
        }
        root = ET.fromstring(xml_text)
        items: list[tuple[str, bool, str]] = []

        for response in root.findall("d:response", ns):
            href_elem = response.find("d:href", ns)
            if href_elem is None or not href_elem.text:
                continue

            resource_type = response.find("d:propstat/d:prop/d:resourcetype", ns)
            is_dir = resource_type is not None and resource_type.find("d:collection", ns) is not None
            
            last_mod_elem = response.find("d:propstat/d:prop/d:getlastmodified", ns)
            last_modified = last_mod_elem.text if last_mod_elem is not None else ""
            
            items.append((href_elem.text, is_dir, last_modified))

        return items
    def _list_files_recursive(self, webdav_root: str, share_token: str, rel_path: str = "") -> list[tuple[str, str]]:
        rel_path_clean = rel_path.strip("/")
        if rel_path_clean:
            encoded_rel_path = "/".join(quote(part) for part in rel_path_clean.split("/"))
            current_url = f"{webdav_root}{encoded_rel_path}/"
        else:
            current_url = webdav_root

        resp = self.session.request(
            method="PROPFIND",
            url=current_url,
            headers={"Depth": "1"},
            auth=(share_token, ""),
            timeout=self.timeout,
        )
        resp.raise_for_status()

        entries = self._propfind_parse(resp.text)
        file_paths: list[tuple[str, str]] = []
        root_prefix = urlparse(webdav_root).path.rstrip("/") + "/"

        for href, is_dir, last_modified in entries[1:]:
            href_path = urlparse(href).path
            if root_prefix not in href_path:
                continue

            relative_item = href_path.split(root_prefix, 1)[1].strip("/")
            if not relative_item:
                continue

            if is_dir:
                file_paths.extend(self._list_files_recursive(webdav_root, share_token, relative_item))
            else:
                file_paths.append((relative_item, last_modified))

        return file_paths

    def _download_webdav_file(self, webdav_root: str, share_token: str, remote_file_path: str) -> requests.Response:
        encoded_remote_path = "/".join(quote(part) for part in remote_file_path.split("/"))
        url = f"{webdav_root}{encoded_remote_path}"
        resp = self.session.get(
            url,
            stream=True,
            auth=(share_token, ""),
            timeout=self.timeout,
            allow_redirects=True,
        )
        resp.raise_for_status()
        return resp

    def download_all_as_pdf(
        self,
        share_url: str,
        output_dir: str = "data/downloads",
        skip_existing: bool = False,
    ) -> list[Path]:
        """Télécharge un partage UTBox, conserve son arborescence et convertit tout en PDF."""
        output_root = Path(output_dir)
        output_root.mkdir(parents=True, exist_ok=True)

        downloaded_files: list[Path] = []

        print(f"Téléchargement depuis UTBox: {share_url}")
        with self.session.get(share_url, stream=True, timeout=self.timeout, allow_redirects=True) as resp:
            resp.raise_for_status()
            archive_name = self._filename_from_headers(resp.headers) or self._filename_from_url(resp.url)
            archive_path = output_root / archive_name

            with open(archive_path, "wb") as f:
                for chunk in resp.iter_content(chunk_size=8 * 1024 * 1024):
                    if chunk:
                        f.write(chunk)

        print("Extraction de l'archive et organisation...")
        if zipfile.is_zipfile(archive_path):
            with zipfile.ZipFile(archive_path, "r") as zf:
                for member in zf.infolist():
                    if member.is_dir():
                        continue

                    relative_path = Path(member.filename)
                    if relative_path.is_absolute() or ".." in relative_path.parts:
                        continue
                    
                    target_path = output_root / relative_path
                    target_path.parent.mkdir(parents=True, exist_ok=True)

                    if skip_existing and target_path.exists():
                        continue

                    if target_path.exists():
                        target_path = self._unique_path(target_path)

                    with zf.open(member, "r") as src, open(target_path, "wb") as dst:
                        self._copy_stream(src, dst)

                    downloaded_files.append(target_path)

            archive_path.unlink(missing_ok=True)
        else:
            filename = archive_path.name
            target_path = self._unique_path(output_root / filename)
            archive_path.replace(target_path)
            downloaded_files.append(target_path)

        print(f"{len(downloaded_files)} fichiers extraits. Conversion en PDF en cours...")
        pdf_files: list[Path] = []
        for file_path in downloaded_files:
            if file_path.suffix.lower() == ".pdf":
                pdf_files.append(file_path)
                continue
                
            pdf_path = file_path.with_suffix(".pdf")


            converted = False
            ext = file_path.suffix.lower()
            try:
                if ext in IMAGE_EXTS:
                    convert_image_to_pdf(file_path, pdf_path)
                    converted = True
                elif ext in TEXT_EXTS:
                    convert_text_to_pdf(file_path, pdf_path)
                    converted = True
                elif ext in OFFICE_EXTS:
                    converted = convert_office_to_pdf(file_path, pdf_path)
            except Exception as e:
                print(f"[ERREUR] Conversion échouée pour {file_path.name}: {e}")

            if converted and pdf_path.exists():
                file_path.unlink(missing_ok=True)
                pdf_files.append(pdf_path)
    def sync_via_zip(
        self,
        share_url: str,
        output_dir: str = "data/downloads"
    ) -> list[Path]:
        """Synchronise de manière incrémentale en utilisant le téléchargement ZIP et la comparaison CRC32."""
        output_root = Path(output_dir)
        output_root.mkdir(parents=True, exist_ok=True)
        
        state_file = output_root / ".utbox_sync_state.json"
        sync_state = {}
        if state_file.exists():
            try:
                sync_state = json.loads(state_file.read_text(encoding="utf-8"))
            except Exception:
                sync_state = {}

        print(f"Téléchargement de l'archive ZIP depuis UTBox: {share_url}")
        
        # S'assurer que le lien finit par /download
        if not share_url.endswith("/download"):
            share_url = share_url.rstrip("/") + "/download"
            
        with self.session.get(share_url, stream=True, timeout=self.timeout, allow_redirects=True) as resp:
            resp.raise_for_status()
            archive_name = self._filename_from_headers(resp.headers) or self._filename_from_url(resp.url)
            archive_path = output_root / archive_name

            with open(archive_path, "wb") as f:
                for chunk in resp.iter_content(chunk_size=8 * 1024 * 1024):
                    if chunk:
                        f.write(chunk)

        print("Analyse de l'archive et détection des nouveautés...")
        new_or_modified_files: list[Path] = []
        new_sync_state = sync_state.copy()

        if zipfile.is_zipfile(archive_path):
            with zipfile.ZipFile(archive_path, "r") as zf:
                for member in zf.infolist():
                    if member.is_dir():
                        continue

                    relative_path = Path(member.filename)
                    if relative_path.is_absolute() or ".." in relative_path.parts:
                        continue
                        
                    rel_path_str = str(relative_path)
                    target_path = output_root / relative_path
                    
                    # Vérification du CRC32 pour détecter les modifications
                    if rel_path_str in sync_state and sync_state[rel_path_str] == member.CRC:
                        continue
                        
                    print(f"Nouveau/Modifié détecté : {rel_path_str}")
                    target_path.parent.mkdir(parents=True, exist_ok=True)

                    with zf.open(member, "r") as src, open(target_path, "wb") as dst:
                        self._copy_stream(src, dst)

                    new_or_modified_files.append(target_path)
                    new_sync_state[rel_path_str] = member.CRC

            archive_path.unlink(missing_ok=True)
        else:
            print("[ERREUR] Le fichier téléchargé n'est pas un ZIP valide.")
            archive_path.unlink(missing_ok=True)
            return []

        if not new_or_modified_files:
            return []

        # Conversion des nouveaux fichiers téléchargés en PDF
        print(f"Conversion de {len(new_or_modified_files)} fichiers...")
        pdf_files: list[Path] = []
        for file_path in new_or_modified_files:
            if file_path.suffix.lower() == ".pdf":
                pdf_files.append(file_path)
                continue
                
            if "RAC" in file_path.parts:
                pdf_files.append(file_path)
                continue
                
            pdf_path = file_path.with_suffix(".pdf")


            converted = False
            ext = file_path.suffix.lower()
            try:
                if ext in IMAGE_EXTS:
                    convert_image_to_pdf(file_path, pdf_path)
                    converted = True
                elif ext in TEXT_EXTS:
                    convert_text_to_pdf(file_path, pdf_path)
                    converted = True
                elif ext in OFFICE_EXTS:
                    converted = convert_office_to_pdf(file_path, pdf_path)
            except Exception as e:
                print(f"[ERREUR] Conversion échouée pour {file_path.name}: {e}")

            if converted and pdf_path.exists():
                file_path.unlink(missing_ok=True)
                pdf_files.append(pdf_path)
            else:
                pdf_files.append(file_path)

        # Sauvegarder l'état
        state_file.write_text(json.dumps(new_sync_state, indent=2, ensure_ascii=False), encoding="utf-8")
        
        return pdf_files

if __name__ == "__main__":
    import json
    project_root = Path(__file__).resolve().parents[1]
    downloads_dir = project_root / "data" / "downloads"

    downloader = UTBoxDownloader(timeout=120)
    files = downloader.sync_webdav_share(
        "https://utbox.univ-tours.fr/s/mmWZHxXHHZSgybT/download",
        output_dir=str(downloads_dir)
    )
    print(f"{len(files)} nouveaux fichiers traités et convertis en PDF dans {downloads_dir}")